"""
Shared fixtures for the course_generator test suite.
Run from the course_generator/ directory:
    pytest                            # skips integration tests
    pytest -m integration             # only integration tests
    pytest -m "not integration"       # only unit/mock tests
"""
import io
import sys
from pathlib import Path

# Make sure the project root is importable regardless of CWD
sys.path.insert(0, str(Path(__file__).parent.parent))

import pytest
from fastapi.testclient import TestClient

from main import app
from services import cache as response_cache

# ---------------------------------------------------------------------------
# Text samples
# ---------------------------------------------------------------------------

ENGLISH_TEXT = (
    "Machine learning is a subset of artificial intelligence that enables systems to learn from data. "
    "Supervised learning uses labeled datasets to train models for classification and regression tasks. "
    "Unsupervised learning discovers hidden patterns in unlabeled data using clustering algorithms. "
    "Reinforcement learning trains agents through a system of rewards and penalties over time. "
    "Deep learning uses neural networks with many layers to process complex patterns efficiently. "
    "Common applications include image recognition, natural language processing, and recommendation systems. "
    "Linear regression predicts continuous values while decision trees handle classification tasks well. "
    "The training process involves minimizing a loss function through gradient descent optimization steps. "
    "Cross-validation helps prevent overfitting by evaluating models on unseen data splits carefully. "
    "Transfer learning allows models pre-trained on large datasets to be fine-tuned for specific tasks."
)

ARABIC_TEXT = (
    "التعلم الآلي هو فرع من فروع الذكاء الاصطناعي يمكن الأنظمة من التعلم من البيانات. "
    "يستخدم التعلم الخاضع للإشراف مجموعات بيانات موسومة لتدريب النماذج على التصنيف والانحدار. "
    "يكتشف التعلم غير الخاضع للإشراف الأنماط المخفية في البيانات غير الموسومة باستخدام التجميع. "
    "يدرب التعلم المعزز العملاء من خلال نظام من المكافآت والعقوبات عبر الزمن. "
    "يستخدم التعلم العميق الشبكات العصبية ذات الطبقات المتعددة لمعالجة الأنماط المعقدة. "
    "تشمل التطبيقات الشائعة التعرف على الصور ومعالجة اللغة الطبيعية وأنظمة التوصية."
)

MIXED_TEXT = (
    "شبكات الحاسوب تستخدم بروتوكول TCP/IP لنقل البيانات عبر الإنترنت بكفاءة عالية. "
    "الراوتر يعمل على توجيه الحزم باستخدام جداول التوجيه والبروتوكولات المختلفة. "
    "يستخدم الـ DHCP لتوزيع عناوين IP تلقائياً على الأجهزة المتصلة بالشبكة المحلية. "
    "الفايروول يحمي الشبكة من الهجمات الخارجية والوصول غير المصرح به للبيانات الحساسة."
)

# ~10 × ENGLISH_TEXT ≈ 1000 words; use 70 repetitions to exceed 6000-word chunk threshold
LONG_TEXT = ("\n\n".join([ENGLISH_TEXT] * 70))


@pytest.fixture
def english_text():
    return ENGLISH_TEXT


@pytest.fixture
def arabic_text():
    return ARABIC_TEXT


@pytest.fixture
def mixed_text():
    return MIXED_TEXT


@pytest.fixture
def long_text():
    return LONG_TEXT


# ---------------------------------------------------------------------------
# File byte fixtures
# ---------------------------------------------------------------------------

@pytest.fixture
def docx_bytes():
    from docx import Document
    doc = Document()
    doc.add_paragraph("Introduction to Machine Learning")
    # Use the shared ENGLISH_TEXT constant so the word count always exceeds MIN_WORDS=50
    doc.add_paragraph(ENGLISH_TEXT)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


@pytest.fixture
def pptx_bytes():
    from pptx import Presentation
    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Machine Learning Overview"
    # Use the shared ENGLISH_TEXT constant so the word count always exceeds MIN_WORDS=50
    slide.placeholders[1].text = ENGLISH_TEXT
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def txt_bytes():
    return ENGLISH_TEXT.encode("utf-8")


@pytest.fixture
def pdf_bytes():
    from fpdf import FPDF
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=12)
    # Use latin-safe text to avoid font issues
    pdf.multi_cell(0, 8, ENGLISH_TEXT)
    return bytes(pdf.output())


# ---------------------------------------------------------------------------
# API client
# ---------------------------------------------------------------------------

@pytest.fixture
def client():
    """A fresh TestClient for each test (lifespan runs but Whisper is cached globally)."""
    with TestClient(app, raise_server_exceptions=False) as c:
        yield c


# ---------------------------------------------------------------------------
# Auto-use fixtures
# ---------------------------------------------------------------------------

@pytest.fixture(autouse=True)
def clear_cache():
    """Wipe the in-memory response cache before and after every test."""
    response_cache._store.clear()
    yield
    response_cache._store.clear()


@pytest.fixture(autouse=True)
def disable_rate_limiting():
    """
    Disable the slowapi rate limiter for every test by default.
    The rate-limit-specific tests must explicitly request `enable_rate_limiting`
    (which depends on this fixture) to re-enable it for the duration of that test.
    """
    from utils.rate_limit import limiter
    original = limiter.enabled
    limiter.enabled = False
    yield
    limiter.enabled = original


@pytest.fixture
def enable_rate_limiting(disable_rate_limiting):
    """
    Re-enable rate limiting and provide a clean counter slate.
    Depends on `disable_rate_limiting` so that fixture's teardown still runs
    and restores the original enabled state afterward.
    """
    from utils.rate_limit import limiter
    limiter._storage.reset()   # wipe any leftover counters from prior tests
    limiter.enabled = True
    yield
