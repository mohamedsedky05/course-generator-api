import io
import re
import unicodedata
from pathlib import Path
from typing import Optional

# Arabic Unicode range
_AR = r'\u0600-\u06ff\u0750-\u077f\u08a0-\u08ff\ufb50-\ufdff\ufe70-\ufeff'


def clean_text(text: str) -> str:
    # NFKC normalises Arabic presentation forms (e.g. ﻻ → لا) and
    # Latin ligatures; also decomposes width-variant ASCII.
    text = unicodedata.normalize("NFKC", text)

    # Normalise line endings
    text = re.sub(r'\r\n|\r', '\n', text)

    # Fix common OCR artifacts: repeated punctuation, stray dashes
    text = re.sub(r'\.{4,}', '...', text)
    text = re.sub(r'-{3,}', '—', text)
    text = re.sub(r'_{2,}', '', text)

    # Ensure a single space between Arabic and Latin characters so mixed
    # sentences like "شبكة192.168" → "شبكة 192.168" are readable
    text = re.sub(rf'([{_AR}])([A-Za-z0-9])', r'\1 \2', text)
    text = re.sub(rf'([A-Za-z0-9])([{_AR}])', r'\1 \2', text)

    # Fix punctuation spacing for LTR text (don't touch RTL-adjacent punctuation)
    text = re.sub(r'([A-Za-z0-9])\s+([,;:!?])', r'\1\2', text)   # "word ," → "word,"
    text = re.sub(r'([,;:!?])(?=[A-Za-z])', r'\1 ', text)          # "word,next" → "word, next"

    # Collapse excessive blank lines and trailing whitespace per line
    lines = [ln.rstrip() for ln in text.split('\n')]
    text = '\n'.join(lines)
    text = re.sub(r'\n{3,}', '\n\n', text)

    # Collapse runs of spaces/tabs (preserve newlines)
    text = re.sub(r'[ \t]+', ' ', text)

    return text.strip()


def extract_from_pdf(content: bytes) -> str:
    import pdfplumber
    text_parts = []
    with pdfplumber.open(io.BytesIO(content)) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
    return clean_text("\n\n".join(text_parts))


def extract_from_docx(content: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(content))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return clean_text("\n\n".join(paragraphs))


def extract_from_pptx(content: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if hasattr(slide, "notes_slide") and slide.notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_texts.append(f"[Notes: {notes}]")
        if slide_texts:
            parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts))
    return clean_text("\n\n".join(parts))


def extract_from_txt(content: bytes) -> str:
    try:
        text = content.decode("utf-8")
    except UnicodeDecodeError:
        text = content.decode("latin-1", errors="replace")
    return clean_text(text)


EXTRACTORS = {
    ".pdf": extract_from_pdf,
    ".docx": extract_from_docx,
    ".pptx": extract_from_pptx,
    ".txt": extract_from_txt,
}

SUPPORTED_EXTENSIONS = set(EXTRACTORS.keys())


def extract_text_from_file(filename: str, content: bytes) -> str:
    ext = Path(filename).suffix.lower()
    if ext not in EXTRACTORS:
        raise ValueError(f"Unsupported file type: {ext}. Supported: {', '.join(SUPPORTED_EXTENSIONS)}")
    return EXTRACTORS[ext](content)


def detect_language(text: str) -> str:
    arabic_chars = sum(1 for c in text if '\u0600' <= c <= '\u06ff')
    total_alpha = sum(1 for c in text if c.isalpha())
    if total_alpha == 0:
        return "en"
    ratio = arabic_chars / total_alpha
    if ratio > 0.4:
        return "ar"
    return "en"


def count_words(text: str) -> int:
    return len(text.split())
